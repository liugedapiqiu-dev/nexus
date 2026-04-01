#!/usr/bin/env python3
"""
PPT Beautifier - PowerPoint 美化和自动化工具
支持字体统一、配色方案、排版优化、模板应用、批量处理等功能
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Optional, List, Dict, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image
except ImportError as e:
    print(f"错误：缺少依赖库 - {e}")
    print("请运行：pip install python-pptx pillow")
    sys.exit(1)


class PPTBeautifier:
    """PPT 美化工具类"""
    
    # 默认配色方案
    DEFAULT_COLORS = {
        'primary': (0, 102, 204),      # 主色 - 蓝色
        'secondary': (0, 170, 136),    # 辅色 - 青色
        'accent': (255, 102, 0),       # 强调色 - 橙色
        'text': (51, 51, 51),          # 文字颜色 - 深灰
        'background': (255, 255, 255), # 背景 - 白色
    }
    
    def __init__(self, input_path: str):
        self.input_path = Path(input_path)
        if not self.input_path.exists():
            raise FileNotFoundError(f"文件不存在：{input_path}")
        
        self.prs = Presentation(self.input_path)
        self.output_path = None
    
    def set_output(self, output_path: str):
        """设置输出路径"""
        self.output_path = Path(output_path)
        return self
    
    def beautify(self, keep_animations: bool = False, keep_transitions: bool = False) -> 'PPTBeautifier':
        """一键美化：统一字体、配色、排版"""
        print("🎨 开始美化 PPT...")
        
        # 统一字体
        self.unify_font(font_name="Microsoft YaHei", title_size=32, body_size=18)
        
        # 应用配色
        self.apply_color_scheme(self.DEFAULT_COLORS)
        
        # 优化排版
        self.optimize_layout()
        
        print("✅ 美化完成")
        return self
    
    def unify_font(self, font_name: str = "Microsoft YaHei", 
                   title_size: int = 32, 
                   body_size: int = 18) -> 'PPTBeautifier':
        """统一字体和字号"""
        print(f"📝 统一字体：{font_name} (标题{title_size}pt, 正文{body_size}pt)")
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = font_name
                            # 根据形状位置判断是标题还是正文
                            if shape.top < Inches(2):
                                run.font.size = Pt(title_size)
                            else:
                                run.font.size = Pt(body_size)
        
        print("✅ 字体统一完成")
        return self
    
    def apply_color_scheme(self, colors: Dict[str, tuple]) -> 'PPTBeautifier':
        """应用配色方案"""
        print(f"🎨 应用配色方案...")
        
        primary = colors.get('primary', self.DEFAULT_COLORS['primary'])
        accent = colors.get('accent', self.DEFAULT_COLORS['accent'])
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                # 形状填充色
                if hasattr(shape, "fill") and shape.fill:
                    try:
                        if shape.shape_type == MSO_SHAPE.RECTANGLE:
                            # 标题形状用主色
                            if shape.top < Inches(1.5):
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = RGBColor(*primary)
                    except:
                        pass
                
                # 文本颜色
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 标题用主色，正文用深色
                            if shape.top < Inches(2):
                                run.font.color.rgb = RGBColor(*primary)
                            else:
                                run.font.color.rgb = RGBColor(*colors.get('text', self.DEFAULT_COLORS['text']))
        
        print("✅ 配色应用完成")
        return self
    
    def optimize_layout(self, align: str = "center") -> 'PPTBeautifier':
        """优化排版布局"""
        print(f"📐 优化排版（对齐：{align}）")
        
        for slide in self.prs.slides:
            # 调整标题位置
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    # 设置文本对齐
                    for paragraph in shape.text_frame.paragraphs:
                        if align == "center":
                            paragraph.alignment = PP_ALIGN.CENTER
                        elif align == "left":
                            paragraph.alignment = PP_ALIGN.LEFT
                        elif align == "right":
                            paragraph.alignment = PP_ALIGN.RIGHT
                    
                    # 调整行距
                    shape.text_frame.paragraphs[0].space_after = Pt(12)
        
        print("✅ 排版优化完成")
        return self
    
    def set_font(self, font_name: str, size: int) -> 'PPTBeautifier':
        """设置统一字体和字号"""
        print(f"📝 设置字体：{font_name} {size}pt")
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = font_name
                            run.font.size = Pt(size)
        
        print("✅ 字体设置完成")
        return self
    
    def resize_images(self, width: int, height: int) -> 'PPTBeautifier':
        """批量调整图片尺寸"""
        print(f"🖼️ 调整图片尺寸：{width}x{height}")
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE.PICTURE:
                    try:
                        shape.width = Inches(width / 96)
                        shape.height = Inches(height / 96)
                    except:
                        pass
        
        print("✅ 图片调整完成")
        return self
    
    def apply_template(self, template_path: str) -> 'PPTBeautifier':
        """应用模板"""
        print(f"📋 应用模板：{template_path}")
        
        template_prs = Presentation(template_path)
        
        # 复制模板的母版布局
        if len(template_prs.slide_layouts) > 0:
            template_layout = template_prs.slide_layouts[0]
            
            # 为每张幻灯片应用模板布局
            for i, slide in enumerate(self.prs.slides):
                try:
                    slide.slide_layout = template_layout
                except:
                    pass
        
        print("✅ 模板应用完成")
        return self
    
    def extract_content(self, output_json: str) -> 'PPTBeautifier':
        """提取幻灯片内容到 JSON"""
        print("📥 提取内容...")
        
        content = []
        for i, slide in enumerate(self.prs.slides):
            slide_data = {
                'slide_number': i + 1,
                'title': '',
                'content': [],
                'shapes_count': len(slide.shapes)
            }
            
            # 提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    if shape.top < Inches(2):
                        slide_data['title'] = shape.text_frame.text.strip()
                    else:
                        slide_data['content'].append(shape.text_frame.text.strip())
            
            content.append(slide_data)
        
        output_path = Path(output_json)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(content, f, ensure_ascii=False, indent=2)
        
        print(f"✅ 内容已提取到：{output_json}")
        return self
    
    def export(self, format: str = 'pdf', output_folder: str = None) -> 'PPTBeautifier':
        """导出为其他格式"""
        print(f"📤 导出为 {format} 格式...")
        
        # python-pptx 本身不支持导出 PDF，这里提供扩展点
        # 实际使用时可以调用 LibreOffice 或其他工具
        if format == 'pdf':
            print("⚠️  注意：导出 PDF 需要 LibreOffice 或 Microsoft PowerPoint")
            print(f"   建议命令：libreoffice --headless --convert-to pdf {self.input_path}")
        elif format in ['png', 'jpg']:
            print("⚠️  注意：导出图片需要额外的转换工具")
        else:
            print(f"❌ 不支持的格式：{format}")
        
        return self
    
    def save(self) -> str:
        """保存文件"""
        if not self.output_path:
            # 生成默认输出路径
            self.output_path = self.input_path.parent / f"beautified_{self.input_path.name}"
        
        self.prs.save(str(self.output_path))
        print(f"💾 已保存：{self.output_path}")
        return str(self.output_path)


def process_batch(input_folder: str, output_folder: str, **kwargs):
    """批量处理文件夹中的 PPT 文件"""
    input_path = Path(input_folder)
    output_path = Path(output_folder)
    output_path.mkdir(parents=True, exist_ok=True)
    
    pptx_files = list(input_path.glob("*.pptx"))
    
    if not pptx_files:
        print(f"❌ 文件夹中没有 PPTX 文件：{input_folder}")
        return
    
    print(f"📁 找到 {len(pptx_files)} 个 PPTX 文件")
    
    for pptx_file in pptx_files:
        try:
            print(f"\n处理：{pptx_file.name}")
            beautifier = PPTBeautifier(str(pptx_file))
            beautifier.set_output(str(output_path / f"beautified_{pptx_file.name}"))
            beautifier.beautify()
            beautifier.save()
        except Exception as e:
            print(f"❌ 处理失败 {pptx_file.name}: {e}")
    
    print(f"\n✅ 批量处理完成！输出目录：{output_folder}")


def main():
    parser = argparse.ArgumentParser(description='PPT 美化工具')
    subparsers = parser.add_subparsers(dest='command', help='命令')
    
    # beautify 命令
    beautify_parser = subparsers.add_parser('beautify', help='一键美化 PPT')
    beautify_parser.add_argument('input', help='输入文件')
    beautify_parser.add_argument('--output', '-o', help='输出文件')
    beautify_parser.add_argument('--keep-animations', action='store_true', help='保留动画')
    beautify_parser.add_argument('--keep-transitions', action='store_true', help='保留过渡效果')
    
    # font 命令
    font_parser = subparsers.add_parser('font', help='统一字体')
    font_parser.add_argument('input', help='输入文件')
    font_parser.add_argument('--font', '-f', default='Microsoft YaHei', help='字体名称')
    font_parser.add_argument('--size', '-s', type=int, default=18, help='字号')
    font_parser.add_argument('--output', '-o', help='输出文件')
    
    # color 命令
    color_parser = subparsers.add_parser('color', help='应用配色方案')
    color_parser.add_argument('input', help='输入文件')
    color_parser.add_argument('--primary', help='主色 (十六进制)')
    color_parser.add_argument('--secondary', help='辅色 (十六进制)')
    color_parser.add_argument('--accent', help='强调色 (十六进制)')
    color_parser.add_argument('--output', '-o', help='输出文件')
    
    # layout 命令
    layout_parser = subparsers.add_parser('layout', help='优化排版')
    layout_parser.add_argument('input', help='输入文件')
    layout_parser.add_argument('--align', '-a', default='center', choices=['left', 'center', 'right'])
    layout_parser.add_argument('--output', '-o', help='输出文件')
    
    # template 命令
    template_parser = subparsers.add_parser('template', help='应用模板')
    template_parser.add_argument('input', help='输入文件')
    template_parser.add_argument('--template', '-t', required=True, help='模板文件')
    template_parser.add_argument('--output', '-o', help='输出文件')
    
    # image 命令
    image_parser = subparsers.add_parser('image', help='处理图片')
    image_parser.add_argument('input', help='输入文件')
    image_parser.add_argument('--resize', '-r', help='调整尺寸 (如：1920x1080)')
    image_parser.add_argument('--output', '-o', help='输出文件')
    
    # batch 命令
    batch_parser = subparsers.add_parser('batch', help='批量处理')
    batch_parser.add_argument('input_folder', help='输入文件夹')
    batch_parser.add_argument('--output', '-o', required=True, help='输出文件夹')
    
    # extract 命令
    extract_parser = subparsers.add_parser('extract', help='提取内容')
    extract_parser.add_argument('input', help='输入文件')
    extract_parser.add_argument('--to', '-t', default='json', help='导出格式')
    extract_parser.add_argument('--output', '-o', help='输出文件')
    
    # export 命令
    export_parser = subparsers.add_parser('export', help='导出文件')
    export_parser.add_argument('input', help='输入文件')
    export_parser.add_argument('--format', '-f', default='pdf', choices=['pdf', 'png', 'jpg'])
    export_parser.add_argument('--output-folder', help='输出文件夹')
    
    args = parser.parse_args()
    
    if not args.command:
        parser.print_help()
        return
    
    try:
        if args.command == 'beautify':
            beautifier = PPTBeautifier(args.input)
            if args.output:
                beautifier.set_output(args.output)
            beautifier.beautify(keep_animations=args.keep_animations, 
                              keep_transitions=args.keep_transitions)
            beautifier.save()
        
        elif args.command == 'font':
            beautifier = PPTBeautifier(args.input)
            if args.output:
                beautifier.set_output(args.output)
            beautifier.set_font(args.font, args.size)
            beautifier.save()
        
        elif args.command == 'color':
            colors = PPTBeautifier.DEFAULT_COLORS.copy()
            if args.primary:
                hex_color = args.primary.lstrip('#')
                colors['primary'] = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            if args.secondary:
                hex_color = args.secondary.lstrip('#')
                colors['secondary'] = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            if args.accent:
                hex_color = args.accent.lstrip('#')
                colors['accent'] = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            
            beautifier = PPTBeautifier(args.input)
            if args.output:
                beautifier.set_output(args.output)
            beautifier.apply_color_scheme(colors)
            beautifier.save()
        
        elif args.command == 'layout':
            beautifier = PPTBeautifier(args.input)
            if args.output:
                beautifier.set_output(args.output)
            beautifier.optimize_layout(align=args.align)
            beautifier.save()
        
        elif args.command == 'template':
            beautifier = PPTBeautifier(args.input)
            if args.output:
                beautifier.set_output(args.output)
            beautifier.apply_template(args.template)
            beautifier.save()
        
        elif args.command == 'image':
            beautifier = PPTBeautifier(args.input)
            if args.output:
                beautifier.set_output(args.output)
            if args.resize:
                width, height = map(int, args.resize.split('x'))
                beautifier.resize_images(width, height)
            beautifier.save()
        
        elif args.command == 'batch':
            process_batch(args.input_folder, args.output)
        
        elif args.command == 'extract':
            beautifier = PPTBeautifier(args.input)
            output_file = args.output or f"{Path(args.input).stem}.{args.to}"
            beautifier.extract_content(output_file)
        
        elif args.command == 'export':
            beautifier = PPTBeautifier(args.input)
            beautifier.export(format=args.format, output_folder=args.output_folder)
        
        print("\n✅ 操作完成！")
    
    except FileNotFoundError as e:
        print(f"❌ 文件错误：{e}")
        sys.exit(1)
    except Exception as e:
        print(f"❌ 错误：{e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
